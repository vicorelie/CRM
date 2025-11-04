#!/usr/bin/env python3
# -*- coding: utf-8 -*-

import sys
import os
import json
import uuid
import subprocess
import logging
from shutil import which

# Configurer le logging pour déboguer les erreurs
logging.basicConfig(filename='extractContent.log', level=logging.ERROR)

# Google Cloud Vision pour OCR
from google.cloud import vision
from google.oauth2 import service_account

# Détection de la langue
from langdetect import detect_langs

# Pour manipuler les images
from PIL import Image

# DOCX via python-docx
import docx

# PyPDF2 pour extraire le texte des PDF (pour les PDF contenant du texte)
try:
    import PyPDF2
except ImportError:
    PyPDF2 = None

# pypandoc pour DOC et PPT (nécessite pandoc)
try:
    import pypandoc
except ImportError:
    pypandoc = None

# xlrd pour XLS
try:
    import xlrd
except ImportError:
    xlrd = None

# openpyxl pour XLSX
try:
    import openpyxl
except ImportError:
    openpyxl = None

# python-pptx pour PPTX
try:
    from pptx import Presentation
except ImportError:
    Presentation = None

# Wand (ImageMagick binding) pour convertir les PDF scannés
try:
    from wand.image import Image as WandImage
except ImportError:
    WandImage = None

def segment_text(text: str):
    """
    Segmente le texte sur les doubles sauts de ligne.
    Retourne une liste de dictionnaires {page, content, section_id}.
    """
    segments = []
    paragraphs = text.split("\n\n")
    section_id = 1
    for p in paragraphs:
        c = p.strip()
        if c:
            segments.append({"page": 0, "content": c, "section_id": section_id})
            section_id += 1
    return segments

def detect_language(full_text: str):
    """
    Détecte la langue via langdetect, filtrée sur [fr, en, ar, ru, he, iw].
    """
    txt = full_text.strip()
    if not txt:
        return "inconnue"
    try:
        possibilities = detect_langs(txt)
        possibilities.sort(key=lambda x: x.prob, reverse=True)
        allowed = ["fr", "en", "ar", "ru", "he", "iw"]
        for p in possibilities:
            if p.lang in allowed:
                return p.lang
        return possibilities[0].lang if possibilities else "inconnue"
    except Exception as e:
        logging.error(f"Erreur détection langue: {e}")
        return "inconnue"

def extract_text_from_image(image_path: str):
    """
    Effectue l'OCR sur une image (jpg, jpeg, png) via Google Cloud Vision.
    """
    try:
        credentials_path = os.path.join(os.path.dirname(__file__), 'config', 'wanatest-447115-1b409ff83dbb.json')
        if not os.path.isfile(credentials_path):
            raise RuntimeError(f"Fichier de credentials introuvable: {credentials_path}")
        credentials = service_account.Credentials.from_service_account_file(credentials_path)
        client = vision.ImageAnnotatorClient(credentials=credentials)
        with open(image_path, "rb") as f:
            content = f.read()
        image = vision.Image(content=content)
        response = client.text_detection(
            image=image, 
            image_context={"language_hints": ["fr", "en", "ar", "ru", "he"]}
        )
        if response.error.message:
            raise RuntimeError(response.error.message)
        annotations = response.text_annotations
        if not annotations:
            return ""
        return annotations[0].description or ""
    except Exception as e:
        logging.error(f"OCR image {image_path}: {e}")
        raise RuntimeError(f"Erreur OCR sur l'image: {e}")

def extract_pdf_text_pypdf(file_path: str):
    """
    Extraction du texte d'un PDF via PyPDF2 (pour les PDF contenant du texte).
    """
    if not PyPDF2:
        raise RuntimeError("Extraction PDF impossible : PyPDF2 n'est pas installé.")
    text_all = []
    try:
        with open(file_path, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            for page in reader.pages:
                txt = page.extract_text()
                if txt:
                    text_all.append(txt)
    except Exception as e:
        logging.error(f"Lecture PDF {file_path}: {e}")
        raise RuntimeError(f"Erreur lecture PDF : {e}")
    return "\n".join(text_all)

def extract_pdf_with_ocr(file_path: str):
    """
    Pour les PDF scannés (sans texte natif), convertit chaque page en image
    puis effectue l'OCR via Google Cloud Vision.
    """
    if WandImage is None:
        raise RuntimeError("Extraction OCR sur PDF impossible : Wand n'est pas installé.")
    all_pages_text = []
    try:
        with WandImage(filename=file_path, resolution=300) as pdf:
            for page in pdf.sequence:
                with WandImage(page) as single_page:
                    blob = single_page.make_blob('png')
                    tmp_path = f"/tmp/pdf_ocr_{uuid.uuid4()}.png"
                    with open(tmp_path, "wb") as f:
                        f.write(blob)
                    page_text = extract_text_from_image(tmp_path)
                    all_pages_text.append(page_text)
                    os.remove(tmp_path)
    except Exception as e:
        logging.error(f"OCR PDF {file_path}: {e}")
        raise RuntimeError(f"Erreur OCR sur PDF scanné: {e}")
    return "\n".join(all_pages_text)

def extract_docx_content(file_path: str):
    """
    Extraction du texte d'un fichier DOCX via python-docx.
    """
    try:
        document = docx.Document(file_path)
    except Exception as e:
        raise RuntimeError(f"Erreur lecture DOCX : {e}")
    segments = []
    page_num = 1
    for para in document.paragraphs:
        txt = para.text.strip()
        if txt:
            segments.append({"page": page_num, "content": txt})
        page_num += 1
    return segments

def extract_doc_content(file_path: str):
    """
    Extraction d'un fichier DOC.
    Convertit le fichier en DOCX via LibreOffice (soffice) puis extrait le contenu.
    """
    base, _ = os.path.splitext(file_path)
    converted_file = base + ".docx"

    if which("soffice") is not None:
        cmd = f'soffice --headless --convert-to docx --outdir "{os.path.dirname(file_path)}" "{file_path}"'
        result = subprocess.run(cmd, shell=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE, universal_newlines=True)
        if result.returncode != 0 or not os.path.isfile(converted_file):
            logging.error("Sortie de soffice: " + result.stdout + " " + result.stderr)
            raise RuntimeError("La conversion du DOC en DOCX a échoué via soffice.")
        return extract_docx_content(converted_file)
    else:
        raise RuntimeError("Extraction DOC impossible : LibreOffice (soffice) n'est pas disponible. Veuillez convertir manuellement le fichier DOC en DOCX.")

def extract_xls_content(file_path: str):
    """
    Extraction d'un fichier XLS (Excel) via xlrd.
    """
    if not xlrd:
        raise RuntimeError("Extraction XLS impossible : xlrd n'est pas installé.")
    try:
        book = xlrd.open_workbook(file_path)
        all_text = []
        for sheet in book.sheets():
            for row_idx in range(sheet.nrows):
                row_values = sheet.row_values(row_idx)
                row_text = " ".join(str(v) for v in row_values if v)
                if row_text:
                    all_text.append(row_text)
        return "\n".join(all_text)
    except Exception as e:
        raise RuntimeError(f"Erreur extraction XLS : {e}")

def extract_xlsx_content(file_path: str):
    """
    Extraction d'un fichier XLSX (Excel) via openpyxl.
    """
    if not openpyxl:
        raise RuntimeError("Extraction XLSX impossible : openpyxl n'est pas installé.")
    try:
        wb = openpyxl.load_workbook(file_path, data_only=True)
        all_text = []
        for sheet in wb.worksheets:
            for row in sheet.iter_rows(values_only=True):
                row_text = " ".join(str(val) for val in row if val is not None)
                if row_text.strip():
                    all_text.append(row_text)
        return "\n".join(all_text)
    except Exception as e:
        raise RuntimeError(f"Erreur extraction XLSX : {e}")

def extract_pptx_content(file_path: str):
    """
    Extraction d'un fichier PPTX via python-pptx.
    """
    if not Presentation:
        raise RuntimeError("Extraction PPTX impossible : python-pptx n'est pas installé.")
    try:
        prs = Presentation(file_path)
        all_text = []
        for slide in prs.slides:
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text.append(shape.text.strip())
            if slide_text:
                all_text.append("\n".join(slide_text))
        return "\n\n".join(all_text)
    except Exception as e:
        raise RuntimeError(f"Erreur extraction PPTX : {e}")

def extract_ppt_content(file_path: str):
    """
    Extraction d'un fichier PPT via pypandoc (nécessite pandoc).
    """
    if not pypandoc:
        raise RuntimeError("Extraction PPT impossible : pypandoc n'est pas installé.")
    try:
        txt = pypandoc.convert_file(file_path, to='plain', format='ppt')
        return txt.strip()
    except Exception as e:
        raise RuntimeError(f"Erreur extraction PPT : {e}")

def main():
    if len(sys.argv) < 2:
        print(json.dumps({"error": "Usage: python extractContent.py <fichier>"}, ensure_ascii=False))
        sys.exit(1)

    file_path = sys.argv[1]
    if not os.path.isfile(file_path):
        print(json.dumps({"error": f"Fichier introuvable: {file_path}"}, ensure_ascii=False))
        sys.exit(1)

    allowed_ext = ["pdf", "doc", "docx", "jpg", "jpeg", "png", "xls", "xlsx", "ppt", "pptx"]
    ext = os.path.splitext(file_path)[1].lower().replace(".", "")
    if ext not in allowed_ext:
        print(json.dumps({"error": f"Format non supporté. Extensions possibles : {', '.join(allowed_ext)}"}, ensure_ascii=False))
        sys.exit(1)

    segments = []
    try:
        if ext == "pdf":
            raw_text = extract_pdf_text_pypdf(file_path).strip()
            if raw_text:
                segments = segment_text(raw_text)
            else:
                ocr_text = extract_pdf_with_ocr(file_path).strip()
                segments = segment_text(ocr_text)
        elif ext in ["jpg", "jpeg", "png"]:
            text = extract_text_from_image(file_path).strip()
            segments = segment_text(text)
        elif ext == "doc":
            segments = extract_doc_content(file_path)
        elif ext == "docx":
            segments = extract_docx_content(file_path)
        elif ext == "xls":
            text = extract_xls_content(file_path)
            segments = segment_text(text)
        elif ext == "xlsx":
            text = extract_xlsx_content(file_path)
            segments = segment_text(text)
        elif ext == "ppt":
            text = extract_ppt_content(file_path)
            segments = segment_text(text)
        elif ext == "pptx":
            text = extract_pptx_content(file_path)
            segments = segment_text(text)

        full_text = " ".join(s["content"] for s in segments)
        lang = detect_language(full_text)
        result = {"language": lang, "segments": segments}
        print(json.dumps(result, ensure_ascii=False, indent=2))
    except Exception as e:
        logging.error(f"Erreur extraction pour {file_path}: {e}")
        print(json.dumps({"error": str(e)}, ensure_ascii=False))
        sys.exit(1)

if __name__ == "__main__":
    main()
