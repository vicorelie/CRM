#!/usr/bin/env python3
# -*- coding: utf-8 -*-

import sys, os, json, subprocess, uuid, logging
from langdetect import detect_langs

# Office
import docx
try:
    import openpyxl
except ImportError:
    openpyxl = None
try:
    from pptx import Presentation
except ImportError:
    Presentation = None
try:
    import pypandoc
except ImportError:
    pypandoc = None

# PDF & OCR (local, sans cloud)
import fitz          # PyMuPDF
import pytesseract
from PIL import Image

logging.basicConfig(filename='extractContent.log', level=logging.ERROR)

ALLOWED_EXT = {"pdf","doc","docx","jpg","jpeg","png","xls","xlsx","ppt","pptx"}

def segment_text(text: str):
    segs, section_id = [], 1
    for p in text.split("\n\n"):
        c = p.strip()
        if c:
            segs.append({"page": 0, "content": c, "section_id": section_id})
            section_id += 1
    return segs

def detect_language(full_text: str):
    txt = (full_text or "").strip()
    if not txt:
        return "inconnue"
    try:
        poss = detect_langs(txt)
        poss.sort(key=lambda x: x.prob, reverse=True)
        allowed = {"fr","en","ar","ru","he","iw"}
        for p in poss:
            if p.lang in allowed:
                return p.lang
        return poss[0].lang if poss else "inconnue"
    except Exception as e:
        logging.error(f"Erreur détection langue: {e}")
        return "inconnue"

# -------- PDF --------

def extract_pdf_text_pymupdf(path: str) -> str:
    out = []
    try:
        with fitz.open(path) as doc:
            for page in doc:
                out.append(page.get_text("text"))
    except Exception as e:
        logging.error(f"PyMuPDF lecture PDF {path}: {e}")
        raise RuntimeError(f"Erreur lecture PDF : {e}")
    return "\n".join(out).strip()

def ocr_pdf_with_ocrmypdf(path: str) -> str:
    base, _ = os.path.splitext(path)
    ocr_path = base + "_ocr.pdf"
    langs = "eng+fra+heb+ara+rus"
    cmd = ["ocrmypdf", "--force-ocr", f"--languages={langs}", "--output-type", "pdf",
           "--skip-text", "--optimize", "3", path, ocr_path]
    r = subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True)
    if r.returncode != 0 or not os.path.isfile(ocr_path):
        logging.error(f"OCRmyPDF échec: code={r.returncode}\nSTDOUT:\n{r.stdout}\nSTDERR:\n{r.stderr}")
        raise RuntimeError("OCR sur PDF scanné: ocrmypdf a échoué.")
    return extract_pdf_text_pymupdf(ocr_path)

# -------- Images --------

def extract_text_from_image(image_path: str) -> str:
    try:
        img = Image.open(image_path)
        return pytesseract.image_to_string(img, lang="eng+fra+heb+ara+rus").strip()
    except Exception as e:
        logging.error(f"OCR image {image_path}: {e}")
        raise RuntimeError(f"Erreur OCR image : {e}")

# -------- Office --------

def extract_docx_content(file_path: str):
    try:
        document = docx.Document(file_path)
    except Exception as e:
        raise RuntimeError(f"Erreur lecture DOCX : {e}")
    segs = []
    page_num = 1
    for para in document.paragraphs:
        txt = (para.text or "").strip()
        if txt:
            segs.append({"page": page_num, "content": txt})
        page_num += 1
    return segs

def soffice_convert_to_txt(src_path: str) -> str:
    outdir = os.path.dirname(src_path)
    cmd = ["soffice", "--headless", "--convert-to", "txt:Text", "--outdir", outdir, src_path]
    r = subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True)
    txt_path = os.path.splitext(src_path)[0] + ".txt"
    if r.returncode != 0 or not os.path.isfile(txt_path):
        logging.error(f"soffice -> txt: code={r.returncode}\nSTDOUT:\n{r.stdout}\nSTDERR:\n{r.stderr}")
        raise RuntimeError("Conversion Office vers TXT a échoué.")
    with open(txt_path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()

def extract_xlsx_text(file_path: str) -> str:
    if not openpyxl:
        raise RuntimeError("openpyxl non installé.")
    try:
        wb = openpyxl.load_workbook(file_path, data_only=True)
        out = []
        for sheet in wb.worksheets:
            for row in sheet.iter_rows(values_only=True):
                row_text = " ".join(str(v) for v in row if v is not None).strip()
                if row_text:
                    out.append(row_text)
        return "\n".join(out)
    except Exception as e:
        raise RuntimeError(f"Erreur extraction XLSX : {e}")

def extract_pptx_text(file_path: str) -> str:
    if not Presentation:
        raise RuntimeError("python-pptx non installé.")
    try:
        prs = Presentation(file_path)
        out = []
        for slide in prs.slides:
            buf = []
            for sh in slide.shapes:
                if hasattr(sh, "text") and sh.text:
                    buf.append(sh.text.strip())
            if buf:
                out.append("\n".join(buf))
        return "\n\n".join(out)
    except Exception as e:
        raise RuntimeError(f"Erreur extraction PPTX : {e}")

def extract_ppt_text_pandoc(file_path: str) -> str:
    if not pypandoc:
        raise RuntimeError("pypandoc non installé.")
    try:
        return pypandoc.convert_file(file_path, to="plain").strip()
    except Exception as e:
        raise RuntimeError(f"Erreur extraction PPT : {e}")

# -------- Main --------

def main():
    if len(sys.argv) < 2:
        print(json.dumps({"error": "Usage: python extractContent.py <fichier>"}, ensure_ascii=False))
        sys.exit(1)

    file_path = sys.argv[1]
    if not os.path.isfile(file_path):
        print(json.dumps({"error": f"Fichier introuvable: {file_path}"}, ensure_ascii=False))
        sys.exit(1)

    ext = os.path.splitext(file_path)[1].lower().replace(".", "")
    if ext not in ALLOWED_EXT:
        print(json.dumps({"error": f"Format non supporté. Extensions possibles : {', '.join(sorted(ALLOWED_EXT))}"}, ensure_ascii=False))
        sys.exit(1)

    try:
        segments = []
        if ext == "pdf":
            text = extract_pdf_text_pymupdf(file_path)
            if not text.strip():
                text = ocr_pdf_with_ocrmypdf(file_path)
            segments = segment_text(text)

        elif ext in {"jpg","jpeg","png"}:
            text = extract_text_from_image(file_path)
            segments = segment_text(text)

        elif ext == "docx":
            segments = extract_docx_content(file_path)

        elif ext == "doc":
            # conversion universelle vers TXT (plus robuste que DOC->DOCX dans pas mal de cas)
            text = soffice_convert_to_txt(file_path)
            segments = segment_text(text)

        elif ext == "xlsx":
            text = extract_xlsx_text(file_path)
            segments = segment_text(text)

        elif ext == "xls":
            # soit xlrd 1.2.0 (si installé), soit conversion vers CSV/TXT
            try:
                text = soffice_convert_to_txt(file_path)
            except Exception:
                # fallback pandoc si dispo
                if pypandoc:
                    text = pypandoc.convert_file(file_path, to="plain")
                else:
                    raise
            segments = segment_text(text)

        elif ext == "pptx":
            # python-pptx ou conversion pandoc
            try:
                text = extract_pptx_text(file_path)
            except Exception:
                if pypandoc:
                    text = pypandoc.convert_file(file_path, to="plain")
                else:
                    raise
            segments = segment_text(text)

        elif ext == "ppt":
            # pandoc recommandé
            text = extract_ppt_text_pandoc(file_path)
            segments = segment_text(text)

        full_text = " ".join(s.get("content","") for s in segments)
        lang = detect_language(full_text)
        print(json.dumps({"language": lang, "segments": segments}, ensure_ascii=False, indent=2))

    except Exception as e:
        logging.error(f"Erreur extraction pour {file_path}: {e}")
        print(json.dumps({"error": str(e)}, ensure_ascii=False))
        sys.exit(1)

if __name__ == "__main__":
    main()
